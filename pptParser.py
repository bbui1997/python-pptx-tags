# TODO:
# 1) Finish this program
#   a. Figure out how we want the user to choose the PowerPoint file
#   b. Error handling involved with choosing a file
#   c. finalize as a group what we want to be our metadata
#   d. figure out/edit how the metadata is being inserted to the program
# 2) Set up code delivery in a way where our user only has to do 1 simple installation to use our program
from __future__ import division, unicode_literals
from textblob import TextBlob as tb
import os
import operator
# from Tkinter import *
# import Tkinter, Tkconstants, tkFileDialog
from stopwords import filter_stop
from pptx import Presentation # pip install python-pptx
from tfidfForStringLists import sortedTfIdfLists
# 1) Open dialog box to select pptx files (currently doesn't parse metadata for .ppt files)
# 2) For each presentation, parse presentation for text
# 3) Take text and categorize into core properties
# 4) Actually insert those core properties to each Presentation
DEFAULT_VAL = 10
def main():
    # Get all of files that user wants to change tags for
    # Tuple is a ( [list of Presentations], [list of filenames] )
    prsAndFileNameTuple = findFile()
    if prsAndFileNameTuple == ([],[]):
        print "No pptx presentations found. Exiting program.\n"
        os.system("pause")
        sys.exit()

    else:
        i = 0
        for i in range(len(prsAndFileNameTuple[1])):
            #print "i is " + str(i)
            prs = prsAndFileNameTuple[0][i]
            #print "prs is " + str(prs)
            filename = prsAndFileNameTuple[1][i]


            print "Parsing: " + os.path.basename(filename)

            # Get the largest text from each slide, add each of those word to a list
            wordList = parseText(prs)

            # rank the word list to figure out the order in which we will populate the tags
            metadata = parseMetaData(wordList, DEFAULT_VAL)

            # Insert metadata (Core Properties) to tags
            populateCoreProperties(prs, metadata, filename)
            i += 1
            print

    print("Successfully added tags to all .pptx files in directory")
    os.system("pause")

def findFile():
    pptx_files = [] # list of filenames
    powerPoints = [] # list of Presentations

    # pptx_files = tkFileDialog.askopenfilenames(filetypes = (("Power Point", "*.pptx"),)) # add each file name to list
    #
    # count = 0
    # for fileName in pptx_files:
    #     print "Selected files " + fileName
    #     powerPoints.append(Presentation(fileName))
    #     count+=1
    #
    # if count == 0:
    #     print("No files entered. Aborting")
    #     return None
    print "Finding all .pptx files.."
    currentPath = os.path.dirname(os.path.abspath(__file__))
    for file in os.listdir(currentPath):
        try:
          if file.endswith(".pptx"):
              pptx_fullPath = os.path.join(currentPath, file)
              pptx_files.append(pptx_fullPath)
              powerPoints.append(Presentation(pptx_fullPath))
        except Exception:
          print "Error with File:" + pptx_fullPath

    # return 2-tuple of lists of Presentations, and filenames
    return (powerPoints, pptx_files)


# for each slide, for each paragraph run, find the runs with the biggest font
# insert each word to a word list
# filter out common words such as "the" or "is" and return that word list
def parseText(presentation):
    splitText = [] # get rid of extra white space
    text = ""
    listOfTexts = []
    greatestFontDict = {}
    for slide in presentation.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    # read the text in the text box, encode it from unicode to ascii, get rid of extra white space
                    # then finally add all the contents in the split() list to the word list
                    # could have data loss
                    splitText = filter_stop(run.text.encode('ascii', 'ignore').split())
                    text += " ".join(splitText)
        text = text.lower()
        listOfTexts.append(tb(text))
    return listOfTexts

# find the X most common words in the word list
# return a String
def parseMetaData(wordList, val):
    sorted_words = sortedTfIdfLists(wordList)
    # initialize loop variables
    keywords = ""

    # This if statement is here to get a string with no extra white space
    if len(sorted_words) != 0:
        keywords = sorted_words[0][0]


    # top 10 or when we run out
    for word in sorted_words[1:val]:
        keywords += ", " + word[0]

    print "Adding tags: " + keywords

    return keywords


# Add all keywords in wordList to tags property and save
# We may need to change this depending on if we want to add any more metadata
# or how we want to insert it
def populateCoreProperties(presentation, metadata, filename):
    string = ""

    ## Useless, right now it's a single string anyways
    # for item in metadata:
    #     string += item
    # presentation.core_properties.keywords = string
    presentation.core_properties.keywords = metadata
    presentation.save(filename)

def printAllKeywords():
    prsAndFileNameTuple = findFile()
    for i in range(len(prsAndFileNameTuple[1])):
        prs = prsAndFileNameTuple[0][i]
        filename = prsAndFileNameTuple[1][i]
        print "hello"

if __name__ == "__main__":
    main()
