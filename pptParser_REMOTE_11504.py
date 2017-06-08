# TODO:
# 1) Finish this program
#   a. Figure out how we want the user to choose the PowerPoint file
#   b. Error handling involved with choosing a file
#   c. finalize as a group what we want to be our metadata
#   d. figure out/edit how the metadata is being inserted to the program
# 2) Set up code delivery in a way where our user only has to do 1 simple installation to use our program

import codecs
import os
import operator
from Tkinter import *
import Tkinter, Tkconstants, tkFileDialog
from stopwords import filter_stop
from pptx import Presentation # pip install python-pptx

# 1) Open dialog box to select pptx files (currently doesn't parse metadata for .ppt files)
# 2) For each presentation, parse presentation for text
# 3) Take text and categorize into core properties
# 4) Actually insert those core properties to each Presentation

def main():
    # Get all of files that user wants to change tags for
    # Tuple is a ( [list of Presentations], [list of filenames] )
    prsAndFileNameTuple = findFile()
    i = 0

    # for each file, parse metadata
    for i in range(len(prsAndFileNameTuple)):
        prs = prsAndFileNameTuple[0][i]
        filename = prsAndFileNameTuple[1][i]

        # Get the largest text from each slide, add each of those word to a list
        wordList = parseText(prs)

        # rank the word list to figure out the order in which we will populate the tags
        metadata = parseMetaData(wordList)

        # Insert metadata (Core Properties) to tags
        populateCoreProperties(prs, metadata, filename)

def findFile():
    pptx_files = [] # list of filenames
    powerPoints = [] # list of Presentations

    pptx_files.extend(tkFileDialog.askopenfilenames()) # add each file name to list

    # add a Presentation object for each filename
    for fileName in pptx_files:
        powerPoints.append(Presentation(fileName))

    # return 2-tuple of lists of Presentations, and filenames
    return (powerPoints, pptx_files)


# for each slide, for each paragraph run, find the runs with the biggest font
# insert each word to a word list
# filter out common words such as "the" or "is" and return that word list
def parseText(presentation):
    wordList = []
    greatestFontDict = {}
    for slide in presentation.slides:
        max = 0
        greatestRun = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:

                for run in paragraph.runs:
                    # read the text in the text box, encode it from unicode to ascii, get rid of extra white space
                    # then finally add all the contents in the split() list to the word list
                    # could have data loss
                    if(run.font.size >= max):
                        greatestRun = run.text.encode('ascii', 'ignore').split()
                        max = run.font.size

        wordList.extend(greatestRun)
    return filter_stop(wordList)

# find the most common words in the word list
# return a String
def parseMetaData(wordList):
    # get the word count
    word_count = {}
    for word in wordList:
        if word in word_count:
            word_count[word] += 1
        else:
            word_count[word] = 1
    sorted_words = sorted(word_count, key=word_count.__getitem__, reverse=True)

    # initialize loop variables
    keywords = ""
    i = 0

    # This if statement is here to get a string with no extra white space
    if len(sorted_words) != 0:
        keywords = sorted_words[0]
        i += 1
    # top 15 or when we run out
    while i < 15 and i < len(sorted_words):
        keywords += ", " + sorted_words[i]
        i += 1

    return keywords


# Add all keywords in wordList to tags property and save
# We may need to change this depending on if we want to add any more metadata
# or how we want to insert it
def populateCoreProperties(presentation, metadata, filename):
    string = ""
    for item in metadata:
        string += item
    presentation.core_properties.keywords = string
    presentation.save(filename)

if __name__ == "__main__":
    main()
